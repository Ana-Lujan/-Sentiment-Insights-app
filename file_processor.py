import os
import tempfile
from typing import Optional

# Try to import optional dependencies
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("PyPDF2 not available. PDF processing disabled.")

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("OCR dependencies not available. Image processing disabled.")

try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False
    print("Whisper not available. Audio/video transcription disabled.")

try:
    import moviepy.editor as mp
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False
    print("MoviePy not available. Video processing disabled.")

# New enterprise format dependencies
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("python-docx not available. DOCX processing disabled.")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    print("pandas not available. Excel/CSV processing disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("python-pptx not available. PPTX processing disabled.")

# Initialize Whisper model for audio/video transcription
whisper_model = None

def get_whisper_model():
    global whisper_model
    if whisper_model is None and WHISPER_AVAILABLE:
        print("Loading Whisper model for audio/video transcription...")
        whisper_model = whisper.load_model("base")
        print("Whisper model loaded.")
    return whisper_model

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file with UTF-8 encoding and error handling."""
    if not PDF_AVAILABLE:
        raise ValueError("PDF processing not available. Install PyPDF2.")
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                # Ensure UTF-8 compatibility and filter problematic characters
                page_text = page_text.encode('utf-8', errors='replace').decode('utf-8')
                # Remove null bytes and other problematic characters
                page_text = page_text.replace('\x00', '').replace('\ufeff', '')
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        raise ValueError(f"Error: Archivo PDF corrupto o cifrado. No se pudo extraer el texto digital. Detalle: {str(e)}")

def extract_text_from_image(file_path: str) -> str:
    """Extract text from image using OCR."""
    if not OCR_AVAILABLE:
        raise ValueError("OCR not available. Install pytesseract and Pillow.")
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang='spa+eng')  # Spanish and English
        return text.strip()
    except Exception as e:
        raise ValueError(f"Error processing image: {str(e)}")

def extract_text_from_audio(file_path: str) -> str:
    """Extract text from audio file using Whisper."""
    if not WHISPER_AVAILABLE:
        raise ValueError("Audio transcription not available. Install openai-whisper.")
    try:
        model = get_whisper_model()
        if model is None:
            raise ValueError("Whisper model not loaded.")
        result = model.transcribe(file_path, language='es')  # Spanish
        return result["text"].strip()
    except Exception as e:
        raise ValueError(f"Error processing audio: {str(e)}")

def extract_text_from_video(file_path: str) -> str:
    """Extract text from video file by extracting audio and transcribing."""
    if not MOVIEPY_AVAILABLE or not WHISPER_AVAILABLE:
        raise ValueError("Video processing not available. Install moviepy and openai-whisper.")
    try:
        # Extract audio from video
        video = mp.VideoFileClip(file_path)
        audio_path = tempfile.mktemp(suffix='.wav')
        video.audio.write_audiofile(audio_path, verbose=False, logger=None)

        # Transcribe audio
        text = extract_text_from_audio(audio_path)

        # Clean up
        os.remove(audio_path)

        return text
    except Exception as e:
        raise ValueError(f"Error processing video: {str(e)}")

def process_file(file_content: bytes, filename: str) -> str:
    """Process uploaded file and extract text based on file type."""
    # Create temporary file
    file_extension = filename.lower().split('.')[-1] if '.' in filename else ''
    with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}' if file_extension else '') as temp_file:
        temp_file.write(file_content)
        temp_file_path = temp_file.name

    try:
        if file_extension == 'pdf':
            return extract_text_from_pdf(temp_file_path)
        elif file_extension == 'docx':
            if DOCX_AVAILABLE:
                return extract_text_from_docx(temp_file_path)
            else:
                raise ValueError("Procesamiento DOCX no disponible. Instale python-docx.")
        elif file_extension == 'xlsx':
            if PANDAS_AVAILABLE:
                return extract_text_from_xlsx(temp_file_path)
            else:
                raise ValueError("Procesamiento XLSX no disponible. Instale pandas y openpyxl.")
        elif file_extension == 'pptx':
            if PPTX_AVAILABLE:
                return extract_text_from_pptx(temp_file_path)
            else:
                raise ValueError("Procesamiento PPTX no disponible. Instale python-pptx.")
        elif file_extension == 'csv':
            if PANDAS_AVAILABLE:
                return extract_text_from_csv(temp_file_path)
            else:
                raise ValueError("Procesamiento CSV no disponible. Instale pandas.")
        elif file_extension in ['jpg', 'jpeg', 'png', 'bmp', 'tiff']:
            if OCR_AVAILABLE:
                text = extract_text_from_image(temp_file_path)
                if len(text.strip()) < 10:
                    raise ValueError("Error OCR: No se detectó texto legible en la imagen. Verifique la calidad o el formato.")
                return text
            else:
                raise ValueError("Contenido multimedia detectado (imagen); requiere procesamiento OCR no disponible.")
        elif file_extension in ['mp3', 'wav', 'flac', 'm4a']:
            if WHISPER_AVAILABLE:
                return extract_text_from_audio(temp_file_path)
            else:
                raise ValueError("Contenido multimedia detectado (audio); requiere procesamiento de voz no disponible.")
        elif file_extension in ['mp4', 'avi', 'mov', 'mkv']:
            if MOVIEPY_AVAILABLE and WHISPER_AVAILABLE:
                return extract_text_from_video(temp_file_path)
            else:
                raise ValueError("Contenido multimedia detectado (video); requiere procesamiento avanzado no disponible.")
        elif file_extension == 'txt':
            # P1 & P2: Force UTF-8 decoding for text files with comprehensive error handling
            try:
                return file_content.decode('utf-8', errors='replace')
            except Exception as e:
                raise ValueError(f"Error: Archivo de texto con codificación incompatible. Use archivos con codificación UTF-8. Detalle: {str(e)}")
        else:
            # For unknown types, try to decode as text
            try:
                text = file_content.decode('utf-8', errors='ignore')
                if text.strip():
                    return text
                else:
                    raise ValueError(f"Archivo no soportado o sin contenido legible: {file_extension}")
            except:
                raise ValueError(f"Archivo no soportado: {file_extension}")
    finally:
        # Clean up temporary file
        try:
            os.remove(temp_file_path)
        except:
            pass

def extract_text_from_csv(file_path: str) -> str:
    """Extract text from CSV file using pandas with UTF-8 encoding."""
    if not PANDAS_AVAILABLE:
        raise ValueError("CSV processing not available. Install pandas.")
    try:
        # P1 & P2: Force UTF-8 encoding for CSV reading
        df = pd.read_csv(file_path, encoding='utf-8')
        # P4: Smart extraction - only text columns, ignore numbers/dates
        text_parts = []
        for col in df.select_dtypes(include=['object', 'string']).columns:
            # Filter out columns that look like they contain only numbers or dates
            sample_values = df[col].dropna().head(5).astype(str)
            if not sample_values.str.match(r'^[\d\.\-\s/]+$').any():  # Skip if looks like numbers/dates
                col_text = df[col].dropna().astype(str).str.cat(sep=' ')
                if col_text.strip():
                    text_parts.append(col_text)
        result = ' '.join(text_parts)
        if not result.strip():
            raise ValueError("Error: No se encontró contenido de texto legible en el archivo CSV.")
        return result
    except UnicodeDecodeError:
        # Fallback for non-UTF-8 files
        try:
            df = pd.read_csv(file_path, encoding='latin-1')
            text_parts = []
            for col in df.select_dtypes(include=['object', 'string']).columns:
                sample_values = df[col].dropna().head(5).astype(str)
                if not sample_values.str.match(r'^[\d\.\-\s/]+$').any():
                    col_text = df[col].dropna().astype(str).str.cat(sep=' ')
                    if col_text.strip():
                        text_parts.append(col_text)
            result = ' '.join(text_parts)
            if not result.strip():
                raise ValueError("Error: No se encontró contenido de texto legible en el archivo CSV.")
            return result
        except Exception as e:
            raise ValueError(f"Error: Archivo CSV corrupto o sin formato válido. No se pudo extraer el texto digital. Detalle: {str(e)}")
    except Exception as e:
        raise ValueError(f"Error: Archivo CSV corrupto o sin formato válido. No se pudo extraer el texto digital. Detalle: {str(e)}")

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file using python-docx."""
    if not DOCX_AVAILABLE:
        raise ValueError("DOCX processing not available. Install python-docx.")
    try:
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text.strip())
        return '\n'.join(text)
    except Exception as e:
        raise ValueError(f"Error: Archivo DOCX corrupto o cifrado. No se pudo extraer el texto digital. Detalle: {str(e)}")

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX file using pandas."""
    if not PANDAS_AVAILABLE:
        raise ValueError("XLSX processing not available. Install pandas and openpyxl.")
    try:
        # P1 & P2: Read Excel file
        df = pd.read_excel(file_path, engine='openpyxl')
        # P4: Smart extraction - only text columns, ignore numbers/dates
        text_parts = []
        for col in df.select_dtypes(include=['object', 'string']).columns:
            # Filter out columns that look like they contain only numbers or dates
            sample_values = df[col].dropna().head(5).astype(str)
            if not sample_values.str.match(r'^[\d\.\-\s/]+$').any():  # Skip if looks like numbers/dates
                col_text = df[col].dropna().astype(str).str.cat(sep=' ')
                if col_text.strip():
                    text_parts.append(col_text)
        result = ' '.join(text_parts)
        if not result.strip():
            raise ValueError("Error: No se encontró contenido de texto legible en el archivo Excel.")
        return result
    except Exception as e:
        raise ValueError(f"Error: Archivo XLSX corrupto o cifrado. No se pudo extraer el texto digital. Detalle: {str(e)}")

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file using python-pptx."""
    if not PPTX_AVAILABLE:
        raise ValueError("PPTX processing not available. Install python-pptx.")
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
        return '\n'.join(text)
    except Exception as e:
        raise ValueError(f"Error: Archivo PPTX corrupto o cifrado. No se pudo extraer el texto digital. Detalle: {str(e)}")